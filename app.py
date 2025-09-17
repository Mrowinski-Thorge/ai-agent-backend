import os
import io
import json
import requests
from flask import Flask, request, send_file, jsonify
from flask_cors import CORS
from groq import Groq
from pptx import Presentation
from pptx.util import Inches

# --- Initialisierung & Konfiguration ---
app = Flask(__name__)
CORS(app, resources={r"/generate": {"origins": "https://mrowinski-thorge.github.io"}})

# --- API-Schlüssel aus Umgebungsvariablen ---
GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
WEBSITE_PASSWORD = os.environ.get("WEBSITE_PASSWORD")
PEXELS_API_KEY = os.environ.get("PEXELS_API_KEY")

if not GROQ_API_KEY or not WEBSITE_PASSWORD:
    raise ValueError("GROQ_API_KEY und WEBSITE_PASSWORD müssen gesetzt sein.")
if not PEXELS_API_KEY:
    print("WARNUNG: PEXELS_API_KEY nicht gefunden. PowerPoint-Bilder werden nicht funktionieren.")

client = Groq(api_key=GROQ_API_KEY)

# --- KORREKTE MODELL-DEFINITIONEN BASIEREND AUF IHRER ANGABE ---
FAST_MODEL = "llama-3.1-8b-instant"
POWERFUL_MODEL = "llama-3.3-70b-versatile" # Korrekter Name für das 70B Modell
VALID_TOOLS = ["retrieval", "code_interpreter"] # Die exakten, gültigen Namen

# --- System-Prompts ---
TRIAGE_SYSTEM_PROMPT = """
Du bist ein Triage-Agent. Deine Aufgabe ist es, eine Benutzeranfrage zu bewerten.
Antworte NUR mit dem Wort 'simple' wenn die Anfrage direkt und ohne Werkzeuge (Websuche, Code-Ausführung), komplexe Planung oder Bild-Erstellung beantwortet werden kann.
Beispiele für 'simple': "Hallo", "Was ist die Hauptstadt von Frankreich?", "Fasse diesen Text zusammen: ...", "Übersetze 'Guten Morgen' ins Englische".
Antworte NUR mit dem Wort 'complex' wenn die Anfrage Werkzeuge, eine PowerPoint-Präsentation, Code-Generierung, eine detaillierte Analyse oder aktuelle Informationen erfordert.
Beispiele für 'complex': "Erstelle eine Präsentation über KI", "Schreibe ein Python-Skript für...", "Wie ist das aktuelle Wetter in Berlin?", "Durchsuche das Web nach...".
Antworte NUR mit 'simple' oder 'complex'.
"""

PLANNER_SYSTEM_PROMPT = f"""
Du bist ein Planungs-Agent. Deine Aufgabe ist es, für eine komplexe Anfrage den besten Ausführungsplan als JSON zu erstellen.
**Verfügbare Werkzeuge:** {VALID_TOOLS}
**JSON-Schema:** {{ "final_tools": ["string"], "optimierter_prompt": "string" }}
**REGELN:**
- Wähle Werkzeuge NUR aus der Liste {VALID_TOOLS}, wenn sie absolut notwendig sind. Sonst gib eine LEERE Liste `[]` zurück.
- Optimiere den Prompt des Benutzers für die nachfolgende Aufgabe.
"""

POWERPOINT_SYSTEM_PROMPT = """
Du bist ein Experte für Präsentationen. Erstelle basierend auf dem User-Prompt ein detailreiches JSON-Objekt. Das JSON muss diesem Schema folgen:
{
  "slides": [
    {
      "title": "Titel der Folie",
      "content": ["Stichpunkt 1", "Stichpunkt 2"],
      "notes": "Sprechernotizen",
      "image_search_query": "kurzer, englischer Suchbegriff für ein Bild, z.B. 'future technology'"
    }
  ]
}
Antworte NUR mit dem validen JSON-Code.
"""

# --- Hilfsfunktionen für Bilder & PowerPoint ---
def search_pexels_image(query):
    if not PEXELS_API_KEY or not query: return None
    try:
        headers = {"Authorization": PEXELS_API_KEY}
        params = {"query": query, "per_page": 1, "orientation": "landscape"}
        response = requests.get("https://api.pexels.com/v1/search", headers=headers, params=params, timeout=10)
        response.raise_for_status()
        data = response.json()
        if data["photos"]:
            image_url = data["photos"][0]["src"]["large"]
            image_response = requests.get(image_url, timeout=10)
            image_response.raise_for_status()
            return io.BytesIO(image_response.content)
    except requests.exceptions.RequestException as e:
        print(f"Pexels API Fehler: {e}")
    return None

def handle_powerpoint_creation(ai_json_response):
    prs = Presentation()
    slides_data = ai_json_response.get('slides', [])
    if not slides_data: return None

    prs.slide_width = Inches(16); prs.slide_height = Inches(9)
    
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = slides_data[0].get('title', 'Präsentation')
    slide.placeholders[1].text = "Erstellt vom Universal AI Agent"

    for slide_info in slides_data:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_info.get('title', '')
        
        tf = slide.placeholders[1].text_frame; tf.clear()
        for point in slide_info.get('content', []):
            p = tf.add_paragraph(); p.text = str(point); p.level = 0
            
        image_stream = search_pexels_image(slide_info.get('image_search_query'))
        if image_stream:
            slide.shapes.add_picture(image_stream, Inches(8), Inches(1.5), height=Inches(4.5))

        slide.notes_slide.notes_text_frame.text = slide_info.get('notes', '')

    ppt_io = io.BytesIO(); prs.save(ppt_io); ppt_io.seek(0)
    return ppt_io

# --- Haupt-Route ---
@app.route('/generate', methods=['POST'])
def generate_agent_response():
    auth_header = request.headers.get('Authorization')
    if not auth_header or f"Bearer {WEBSITE_PASSWORD}" != auth_header:
        return jsonify({"error": "Ungültige Authentifizierung"}), 401

    data = request.get_json()
    user_prompt = data.get('prompt')
    output_format = data.get('output_format', 'text')

    if not user_prompt: return jsonify({"error": "Kein Prompt angegeben"}), 400

    try:
        # --- Triage-Stufe ---
        triage_messages = [{"role": "system", "content": TRIAGE_SYSTEM_PROMPT}, {"role": "user", "content": user_prompt}]
        triage_completion = client.chat.completions.create(model=FAST_MODEL, messages=triage_messages, temperature=0.0)
        decision = triage_completion.choices[0].message.content.strip().lower()

        # --- Pfad 1: Einfache Anfrage ---
        if decision == 'simple' and output_format == 'text':
            simple_messages = [{"role": "system", "content": "Du bist ein hilfreicher Assistent. Antworte direkt und präzise."}, {"role": "user", "content": user_prompt}]
            simple_completion = client.chat.completions.create(model=FAST_MODEL, messages=simple_messages)
            return jsonify({"responseText": simple_completion.choices[0].message.content})

        # --- Pfad 2: Komplexe Anfrage ---
        optimierter_prompt = user_prompt
        final_tools_names = []

        if output_format != 'powerpoint':
            planner_messages = [{"role": "system", "content": PLANNER_SYSTEM_PROMPT}, {"role": "user", "content": f"User-Prompt: \"{user_prompt}\", Output-Format: \"{output_format}\""}]
            planner_completion = client.chat.completions.create(model=FAST_MODEL, messages=planner_messages, response_format={"type": "json_object"})
            plan = json.loads(planner_completion.choices[0].message.content)
            
            optimierter_prompt = plan.get("optimierter_prompt", user_prompt)
            # **ROBUSTE VALIDIERUNG:** Nur gültige Werkzeuge werden übernommen.
            final_tools_names = [tool for tool in plan.get("final_tools", []) if tool in VALID_TOOLS]

        # --- Executor-Phase ---
        system_prompt = "Du bist ein Weltklasse-Experte."
        if output_format == 'powerpoint': system_prompt = POWERPOINT_SYSTEM_PROMPT
        elif output_format == 'code': system_prompt = "Du bist ein erfahrener Software-Entwickler."
        
        executor_messages = [{"role": "system", "content": system_prompt}, {"role": "user", "content": optimierter_prompt}]
        
        completion_params = {"model": POWERFUL_MODEL, "messages": executor_messages}
        if final_tools_names:
            completion_params["tools"] = [{"type": name} for name in final_tools_names]
        
        if output_format == 'powerpoint':
            completion_params["response_format"] = {"type": "json_object"}

        executor_completion = client.chat.completions.create(**completion_params)
        ai_response_content = executor_completion.choices[0].message.content

        # --- Ausgabe-Verarbeitung ---
        if output_format == 'powerpoint':
            ai_json = json.loads(ai_response_content)
            ppt_file = handle_powerpoint_creation(ai_json)
            return send_file(ppt_file, as_attachment=True, download_name='praesentation.pptx', mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        else:
            return jsonify({"responseText": ai_response_content})

    except Exception as e:
        print(f"Ein Fehler ist aufgetreten: {e}")
        return jsonify({"error": "Ein interner Serverfehler ist aufgetreten.", "details": str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5001)
